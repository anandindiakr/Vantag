"""Build Vantag platform overview presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─── Brand palette ──────────────────────────────────────────────────────
VIOLET = RGBColor(0x5B, 0x21, 0xB6)   # SG primary
SAFFRON = RGBColor(0xF5, 0x9E, 0x0B)   # IN accent
EMERALD = RGBColor(0x10, 0xB9, 0x81)   # MY primary
INK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
PAPER = RGBColor(0xF9, 0xFA, 0xFB)
CARD = RGBColor(0xFF, 0xFF, 0xFF)
BORDER = RGBColor(0xE5, 0xE7, 0xEB)
DANGER = RGBColor(0xDC, 0x26, 0x26)
DARK = RGBColor(0x11, 0x18, 0x27)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

# ─── helpers ────────────────────────────────────────────────────────────
def add_slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    bgshape = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bgshape.fill.solid(); bgshape.fill.fore_color.rgb = bg
    bgshape.line.fill.background()
    bgshape.shadow.inherit = False
    return s

def txt(slide, x, y, w, h, text, size=14, color=INK, bold=False, align=PP_ALIGN.LEFT,
        font="Calibri", anchor=MSO_ANCHOR.TOP, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
    return tb

def rect(slide, x, y, w, h, fill=CARD, line=None, rounded=True, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        sh.adjustments[0] = 0.08
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    if not shadow:
        sh.shadow.inherit = False
    return sh

def circle(slide, x, y, d, fill=VIOLET):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def pill(slide, x, y, text, fill=VIOLET, color=CARD, size=9, w=None):
    est_w = w if w else max(0.9, 0.11 * len(text) + 0.2)
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(est_w), Inches(0.32))
    sh.adjustments[0] = 0.5
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.name = "Calibri"; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = color
    return sh

def line(slide, x1, y1, x2, y2, color=MUTED, weight=1.25, dashed=False):
    from pptx.util import Emu
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    if dashed:
        ln.line.dash_style = 7  # dash
    return ln

def footer(slide, num, total):
    # brand footer bar
    rect(slide, 0, 7.2, 13.333, 0.3, fill=DARK, rounded=False)
    txt(slide, 0.4, 7.22, 6, 0.28, "VANTAG  |  Retail Intelligence Platform",
        size=9, color=CARD, bold=True)
    txt(slide, 11.5, 7.22, 1.5, 0.28, f"{num} / {total}",
        size=9, color=CARD, align=PP_ALIGN.RIGHT)

TOTAL = 18
N = [0]
def nextn():
    N[0] += 1; return N[0]

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(bg=DARK)
# gradient stripe
rect(s, 0, 0, 13.333, 0.15, fill=VIOLET, rounded=False)
rect(s, 4.44, 0, 4.44, 0.15, fill=SAFFRON, rounded=False)
rect(s, 8.88, 0, 4.44, 0.15, fill=EMERALD, rounded=False)
# big V mark
circle(s, 0.7, 2.5, 2.5, fill=VIOLET)
txt(s, 0.7, 2.5, 2.5, 2.5, "V", size=120, color=CARD, bold=True,
    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font="Calibri")
# title block
txt(s, 3.6, 2.2, 9, 0.5, "VANTAG", size=14, color=SAFFRON, bold=True)
txt(s, 3.6, 2.7, 9, 1.4, "Retail Intelligence\n& Security Platform",
    size=54, color=CARD, bold=True, font="Calibri")
txt(s, 3.6, 4.8, 9, 0.5,
    "AI-Powered · Edge-First · Multi-Region SaaS",
    size=18, color=SAFFRON, italic=True)
txt(s, 3.6, 5.4, 9, 0.4,
    "Turning any IP camera into a 24/7 AI guardian",
    size=14, color=CARD)
# region pills bottom
pill(s, 3.6, 6.2, "🇮🇳 India — Retail Nazar", fill=SAFFRON, color=DARK, size=11, w=2.9)
pill(s, 6.7, 6.2, "🇸🇬 Singapore — Vantag", fill=VIOLET, color=CARD, size=11, w=2.8)
pill(s, 9.6, 6.2, "🇲🇾 Malaysia — JagaJaga", fill=EMERALD, color=CARD, size=11, w=3.1)
txt(s, 0.4, 7.15, 12.5, 0.3, "Investor & Partner Overview  ·  Version 1.0",
    size=10, color=MUTED, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "THE PROBLEM", size=11, color=VIOLET, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "Retail shrinkage costs trillions.\nTraditional CCTV rarely catches it.",
    size=32, color=INK, bold=True)

# Left stat cards
stats = [
    ("USD 112B", "global retail shrinkage\nannual loss", DANGER),
    ("1.5–3%", "revenue lost every year\nin emerging markets", SAFFRON),
    ("< 3%", "of CCTV footage is\nactually reviewed", VIOLET),
    ("USD 5K–50K", "typical cost per store for\nexisting AI-CCTV systems", INK),
]
for i, (big, sub, color) in enumerate(stats):
    x = 0.6 + i * 3.1
    rect(s, x, 2.3, 2.9, 2.2, fill=CARD, line=BORDER)
    txt(s, x + 0.2, 2.5, 2.5, 1.0, big, size=34, color=color, bold=True)
    txt(s, x + 0.2, 3.55, 2.5, 0.9, sub, size=12, color=MUTED)

# callout
rect(s, 0.6, 4.95, 12.2, 1.8, fill=DARK, line=None)
txt(s, 1.0, 5.1, 11.5, 0.45, "WHO SUFFERS MOST", size=11, color=SAFFRON, bold=True)
txt(s, 1.0, 5.55, 11.5, 1.2,
    "Small retailers, mall managers, hospitals, clinics, post offices, police stations\n"
    "and homes — those who need protection most — are priced out of the market.",
    size=15, color=CARD)
footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "OUR SOLUTION", size=11, color=EMERALD, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "DIY AI security in under 30 minutes.",
    size=32, color=INK, bold=True)
txt(s, 0.6, 1.7, 12, 0.5,
    "Works with any existing IP camera. No installers, no NVR lock-in, no monthly data bills.",
    size=15, color=MUTED)

# Three-pillar row
pillars = [
    ("🛡", "EDGE-FIRST", "All AI runs on-premises.\nVideo never leaves your store.", VIOLET),
    ("🔌", "PLUG & PLAY", "Auto-discovers cameras,\nauto-draws zones, auto-starts.", SAFFRON),
    ("💳", "SUBSCRIPTION", "Razorpay in INR, SGD, MYR.\nFrom ₹1,999 / S$39 / RM59.", EMERALD),
]
for i, (icon, title, desc, color) in enumerate(pillars):
    x = 0.6 + i * 4.1
    rect(s, x, 2.6, 3.9, 3.5, fill=CARD, line=BORDER)
    circle(s, x + 0.3, 2.85, 0.9, fill=color)
    txt(s, x + 0.3, 2.85, 0.9, 0.9, icon, size=32, color=CARD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.3, 4.0, 3.5, 0.5, title, size=17, color=color, bold=True)
    txt(s, x + 0.3, 4.55, 3.5, 1.5, desc, size=13, color=INK)

# bottom stripe
rect(s, 0.6, 6.3, 12.2, 0.7, fill=DARK, rounded=True)
txt(s, 0.6, 6.38, 12.2, 0.55,
    "12 AI Detectors  ·  11 Languages  ·  3 Regions  ·  Android / Windows / Pi / Jetson",
    size=13, color=CARD, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — HOW IT WORKS (ARCHITECTURE)
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "HOW IT WORKS", size=11, color=VIOLET, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "Edge AI + Cloud Dashboard",
    size=32, color=INK, bold=True)

# left: customer premises
rect(s, 0.5, 2.0, 5.5, 4.8, fill=CARD, line=BORDER)
txt(s, 0.7, 2.15, 5, 0.4, "CUSTOMER PREMISES", size=10, color=VIOLET, bold=True)
# cameras
for i, lbl in enumerate(["Cam 1", "Cam 2", "Cam 3", "NVR"]):
    circle(s, 0.9, 2.7 + i * 0.75, 0.5, fill=EMERALD)
    txt(s, 0.9, 2.7 + i * 0.75, 0.5, 0.5, "📹", size=18, color=CARD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, 1.5, 2.78 + i * 0.75, 1.5, 0.35, lbl, size=11, color=INK, bold=True)
# edge agent
rect(s, 3.3, 3.2, 2.5, 2.4, fill=VIOLET, line=None)
txt(s, 3.3, 3.3, 2.5, 0.4, "EDGE AGENT", size=11, color=SAFFRON, bold=True, align=PP_ALIGN.CENTER)
txt(s, 3.3, 3.75, 2.5, 0.4, "Android · Windows", size=10, color=CARD, align=PP_ALIGN.CENTER)
txt(s, 3.3, 4.1, 2.5, 0.4, "Raspberry Pi · Jetson", size=10, color=CARD, align=PP_ALIGN.CENTER)
txt(s, 3.3, 4.5, 2.5, 0.4, "YOLOv8 Inference", size=10, color=CARD, italic=True, align=PP_ALIGN.CENTER)
pill(s, 3.95, 5.1, "RTSP · ONVIF", fill=CARD, color=VIOLET, size=9, w=1.25)
# arrows cameras → agent
for i in range(4):
    line(s, 1.45, 2.95 + i * 0.75, 3.3, 4.3, color=MUTED, weight=1)

# right: cloud
rect(s, 7.5, 2.0, 5.3, 4.8, fill=CARD, line=BORDER)
txt(s, 7.7, 2.15, 5, 0.4, "VANTAG CLOUD", size=10, color=EMERALD, bold=True)
services = [
    ("FastAPI Backend", VIOLET, 2.55),
    ("PostgreSQL + Redis", INK, 3.20),
    ("Mosquitto MQTT", EMERALD, 3.85),
    ("Razorpay Gateway", SAFFRON, 4.50),
    ("React Web + PWA Mobile", VIOLET, 5.15),
    ("Super-Admin Panel", DARK, 5.80),
]
for lbl, c, y in services:
    rect(s, 7.8, y, 4.7, 0.55, fill=PAPER, line=BORDER)
    circle(s, 7.95, y + 0.1, 0.35, fill=c)
    txt(s, 8.4, y + 0.13, 4, 0.35, lbl, size=12, color=INK, bold=True)

# middle arrow: encrypted
line(s, 6.05, 4.4, 7.5, 4.4, color=EMERALD, weight=3)
txt(s, 6.0, 3.9, 1.5, 0.4, "HTTPS\nmetadata only", size=9, color=EMERALD, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, 6.0, 4.6, 1.5, 0.3, "🔒", size=16, color=EMERALD, align=PP_ALIGN.CENTER)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — 12 AI DETECTORS
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "AI CAPABILITIES", size=11, color=VIOLET, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "12 real-time detectors · > 92% accuracy",
    size=32, color=INK, bold=True)

detectors = [
    ("🎯", "Product Sweeping", "Theft by grab-and-run", DANGER),
    ("⏱", "Anomalous Dwell", "Loitering in zones", SAFFRON),
    ("📦", "Empty Shelf", "Out-of-stock alerts", EMERALD),
    ("🦹", "Shoplifting", "Concealment patterns", DANGER),
    ("↔️", "Inventory Move", "Cross-zone transfers", VIOLET),
    ("🤕", "Fall Detection", "Customer / staff falls", DANGER),
    ("🚷", "Zone Intrusion", "Restricted-area entry", SAFFRON),
    ("👥", "Crowd Density", "Queue & occupancy", VIOLET),
    ("📷", "Camera Tamper", "Lens block / cable cut", DANGER),
    ("👔", "Staff Behaviour", "Till away, breaks", EMERALD),
    ("🌙", "After-Hours", "Closed-window motion", DANGER),
    ("🚗", "License Plate", "Vehicle recognition", VIOLET),
]
for i, (icon, title, desc, color) in enumerate(detectors):
    row, col = divmod(i, 4)
    x = 0.6 + col * 3.08
    y = 2.0 + row * 1.55
    rect(s, x, y, 2.95, 1.4, fill=CARD, line=BORDER)
    circle(s, x + 0.15, y + 0.18, 0.65, fill=color)
    txt(s, x + 0.15, y + 0.18, 0.65, 0.65, icon, size=18, color=CARD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.9, y + 0.22, 2.0, 0.4, title, size=12, color=INK, bold=True)
    txt(s, x + 0.9, y + 0.62, 2.0, 0.6, desc, size=10, color=MUTED)

txt(s, 0.6, 6.8, 12, 0.3,
    "Base model: YOLOv8n-seg, fine-tuned on 45K retail frames · BoT-SORT multi-object tracking",
    size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — 30-MIN ONBOARDING
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "CUSTOMER JOURNEY", size=11, color=SAFFRON, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "From signup to live AI in 30 minutes",
    size=32, color=INK, bold=True)

steps = [
    ("01", "REGISTER", "Shop name, email,\nregion auto-detected", "< 3 min"),
    ("02", "PICK PLAN", "Starter · Growth · Pro\nRazorpay checkout", "< 5 min"),
    ("03", "DOWNLOAD", "Edge Agent for\nphone / PC / Pi", "< 5 min"),
    ("04", "SCAN QR", "One-tap pairing\nto cloud", "< 2 min"),
    ("05", "AUTO-DISCOVER", "ONVIF + ARP scan\nfinds cameras", "< 5 min"),
    ("06", "DRAW ZONES", "Click-click polygons\non snapshots", "< 10 min"),
]
for i, (num, title, desc, time) in enumerate(steps):
    row, col = divmod(i, 3)
    x = 0.6 + col * 4.15
    y = 2.1 + row * 2.3
    # big number circle
    circle(s, x, y, 1.2, fill=VIOLET if i % 2 == 0 else SAFFRON)
    txt(s, x, y, 1.2, 1.2, num, size=28, color=CARD, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rect(s, x + 1.4, y, 2.5, 1.9, fill=CARD, line=BORDER)
    txt(s, x + 1.6, y + 0.1, 2.3, 0.4, title, size=13, color=INK, bold=True)
    txt(s, x + 1.6, y + 0.55, 2.3, 0.9, desc, size=11, color=MUTED)
    pill(s, x + 1.6, y + 1.5, time, fill=EMERALD, color=CARD, size=9, w=1.0)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — THREE REGIONS
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(bg=DARK)
txt(s, 0.6, 0.35, 12, 0.5, "GEOGRAPHIC FOOTPRINT", size=11, color=SAFFRON, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "One platform, three localised brands",
    size=32, color=CARD, bold=True)

regions = [
    ("🇮🇳", "RETAIL NAZAR", "retailnazar.com / .in / .info",
     "Hindi · Tamil · Telugu · Kannada\nMalayalam · Marathi · Gujarati\nBengali · Punjabi · English",
     "₹ INR",  "7M shops TAM", SAFFRON),
    ("🇸🇬", "VANTAG", "retail-vantag.com",
     "English",
     "S$ SGD", "200K shops TAM", VIOLET),
    ("🇲🇾", "JAGAJAGA", "jagajaga.my\nretailjagajaga.com",
     "Malay · English",
     "RM MYR", "1M shops TAM", EMERALD),
]
for i, (flag, brand, domain, langs, cur, mkt, color) in enumerate(regions):
    x = 0.6 + i * 4.15
    rect(s, x, 2.0, 3.95, 4.8, fill=INK, line=color)
    txt(s, x + 0.2, 2.15, 3.5, 0.8, flag, size=48,
        align=PP_ALIGN.LEFT)
    txt(s, x + 0.2, 3.0, 3.6, 0.6, brand, size=20, color=color, bold=True)
    txt(s, x + 0.2, 3.55, 3.6, 0.6, domain, size=10, color=MUTED, italic=True)
    rect(s, x + 0.2, 4.15, 3.55, 0.04, fill=color, rounded=False)
    txt(s, x + 0.2, 4.3, 3.55, 0.35, "LANGUAGES", size=9, color=MUTED, bold=True)
    txt(s, x + 0.2, 4.6, 3.55, 1.3, langs, size=11, color=CARD)
    pill(s, x + 0.2, 6.0, cur, fill=color, color=CARD, size=10, w=1.1)
    pill(s, x + 1.4, 6.0, mkt, fill=CARD, color=INK, size=10, w=1.8)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — PRICING
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "PRICING", size=11, color=VIOLET, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "Three tiers. Local currency. No lock-in.",
    size=32, color=INK, bold=True)

# Header row
hdr = [("PLAN", 0.6), ("CAMERAS", 2.6), ("INDIA", 4.2), ("SINGAPORE", 6.7), ("MALAYSIA", 9.2), ("FEATURES", 11.4)]
rect(s, 0.5, 2.0, 12.3, 0.55, fill=DARK, rounded=True)
for label, x in hdr:
    txt(s, x, 2.08, 2.5, 0.4, label, size=11, color=CARD, bold=True)

plans = [
    ("STARTER",  "5",  "₹1,999/mo",  "S$39/mo",  "RM59/mo",  "Core 12 detectors\n30-day history", PAPER),
    ("GROWTH",   "15", "₹4,499/mo",  "S$99/mo",  "RM149/mo", "+ POS integration\n+ Door-lock MQTT\n90-day history", CARD),
    ("PRO",      "30", "₹8,999/mo",  "S$189/mo", "RM299/mo", "+ Priority support\n+ API access\n365-day history", VIOLET),
]
for i, (plan, cams, inr, sgd, myr, feats, bg) in enumerate(plans):
    y = 2.65 + i * 1.48
    text_color = CARD if bg == VIOLET else INK
    rect(s, 0.5, y, 12.3, 1.38, fill=bg, line=BORDER)
    txt(s, 0.7, y + 0.15, 1.8, 0.4, plan, size=18,
        color=VIOLET if bg != VIOLET else SAFFRON, bold=True)
    if i == 2:
        pill(s, 0.7, y + 0.65, "MOST POPULAR", fill=SAFFRON, color=DARK, size=8, w=1.5)
    txt(s, 2.6, y + 0.45, 1.5, 0.5, cams, size=22, color=text_color, bold=True)
    txt(s, 4.2, y + 0.5, 2.4, 0.5, inr, size=18, color=text_color, bold=True)
    txt(s, 6.7, y + 0.5, 2.4, 0.5, sgd, size=18, color=text_color, bold=True)
    txt(s, 9.2, y + 0.5, 2.4, 0.5, myr, size=18, color=text_color, bold=True)
    txt(s, 11.4, y + 0.15, 1.8, 1.15, feats, size=10,
        color=text_color if bg == VIOLET else MUTED)

txt(s, 0.6, 7.0, 12.2, 0.25,
    "Annual plans save 17% · 7-day money-back · Razorpay: UPI, cards, net-banking, FPX, PayNow",
    size=10, color=MUTED, italic=True, align=PP_ALIGN.CENTER)
footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "MARKET OPPORTUNITY", size=11, color=EMERALD, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "8.2 million shops across 3 regions",
    size=32, color=INK, bold=True)

# Left: TAM/SAM/SOM funnel
rect(s, 0.6, 2.1, 6.2, 4.9, fill=CARD, line=BORDER)
txt(s, 0.8, 2.2, 5.8, 0.4, "ADDRESSABLE MARKET FUNNEL", size=10, color=VIOLET, bold=True)
funnel = [
    ("TAM", "8.2M shops",     VIOLET,  5.5, 2.8),
    ("SAM", "85K reachable",  SAFFRON, 4.3, 4.0),
    ("SOM", "18K target Y3",  EMERALD, 3.1, 5.2),
]
for tag, val, c, w, y in funnel:
    cx = 0.6 + 6.2 / 2 - w / 2
    rect(s, cx, y, w, 0.9, fill=c, rounded=True)
    txt(s, cx, y + 0.12, w, 0.4, tag, size=14, color=CARD, bold=True, align=PP_ALIGN.CENTER)
    txt(s, cx, y + 0.48, w, 0.4, val, size=14, color=CARD, align=PP_ALIGN.CENTER)

# Right: per-region table
rect(s, 7.0, 2.1, 5.8, 4.9, fill=CARD, line=BORDER)
txt(s, 7.2, 2.2, 5.4, 0.4, "BY REGION — YEAR 3 TARGETS", size=10, color=EMERALD, bold=True)
rows_data = [
    ("Region", "Shops", "Target", "Revenue"),
    ("India", "7.0M", "15K", "USD 4.8M"),
    ("Malaysia", "1.0M", "2K", "USD 1.0M"),
    ("Singapore", "0.2M", "1K", "USD 1.4M"),
    ("TOTAL", "8.2M", "18K", "USD 7.2M"),
]
col_x = [7.2, 8.5, 9.9, 11.1]
for ri, row in enumerate(rows_data):
    y = 2.75 + ri * 0.7
    bg = DARK if ri == 0 else (PAPER if ri == len(rows_data) - 1 else CARD)
    tc = CARD if ri == 0 else INK
    if ri == 0 or ri == len(rows_data) - 1:
        rect(s, 7.0, y - 0.05, 5.8, 0.7, fill=bg, rounded=False, line=BORDER if ri != 0 else None)
    for ci, val in enumerate(row):
        bold = ri == 0 or ri == len(rows_data) - 1
        txt(s, col_x[ci], y + 0.08, 1.5, 0.5, val, size=12, color=tc, bold=bold)
txt(s, 7.2, 6.35, 5.4, 0.35, "At median ARPU of USD 400 / year", size=10, color=MUTED, italic=True)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "TECHNOLOGY", size=11, color=VIOLET, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "Modern, open, production-grade stack",
    size=32, color=INK, bold=True)

stacks = [
    ("BACKEND",    VIOLET, ["FastAPI · async Python 3.11", "PostgreSQL 15 + asyncpg",
                             "Redis pub/sub + cache", "Mosquitto MQTT 5", "JWT + bcrypt auth"]),
    ("FRONTEND",   SAFFRON, ["React 18 + TypeScript", "Vite build, Tailwind CSS",
                              "PWA installable mobile", "react-i18next (11 langs)", "Chart.js + Mapbox"]),
    ("AI / ML",    EMERALD, ["YOLOv8n-seg (Ultralytics)", "BoT-SORT tracker",
                              "TensorRT on Jetson (7× speed)", "ONNX export for Android",
                              "45K-frame retail dataset"]),
    ("DEVOPS",     INK,     ["Hostinger VPS (Ubuntu 22)", "Nginx + Let's Encrypt TLS",
                              "systemd services", "Prometheus + Sentry", "GitHub CI/CD"]),
]
for i, (title, color, items) in enumerate(stacks):
    row, col = divmod(i, 2)
    x = 0.6 + col * 6.15
    y = 2.1 + row * 2.5
    rect(s, x, y, 5.95, 2.3, fill=CARD, line=BORDER)
    rect(s, x, y, 0.15, 2.3, fill=color, rounded=False)
    txt(s, x + 0.35, y + 0.1, 5.4, 0.45, title, size=15, color=color, bold=True)
    for j, item in enumerate(items):
        circle(s, x + 0.4, y + 0.65 + j * 0.32, 0.1, fill=color)
        txt(s, x + 0.65, y + 0.56 + j * 0.32, 5, 0.35, item, size=11, color=INK)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — SECURITY & PRIVACY
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "SECURITY & PRIVACY", size=11, color=DANGER, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "Video never leaves your premises.",
    size=32, color=INK, bold=True)

# Big diagram
rect(s, 0.6, 2.0, 12.2, 2.3, fill=DARK, rounded=True)
# Customer side
circle(s, 1.2, 2.6, 1.1, fill=EMERALD)
txt(s, 1.2, 2.6, 1.1, 1.1, "🏪", size=30, color=CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 0.8, 3.75, 2, 0.3, "YOUR SHOP", size=10, color=EMERALD, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0.8, 4.0, 2, 0.25, "Video stays here", size=9, color=CARD, align=PP_ALIGN.CENTER)
# Cloud side
circle(s, 11, 2.6, 1.1, fill=VIOLET)
txt(s, 11, 2.6, 1.1, 1.1, "☁", size=36, color=CARD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
txt(s, 10.6, 3.75, 2, 0.3, "VANTAG CLOUD", size=10, color=VIOLET, bold=True, align=PP_ALIGN.CENTER)
txt(s, 10.6, 4.0, 2, 0.25, "Metadata only", size=9, color=CARD, align=PP_ALIGN.CENTER)
# Arrow between
line(s, 2.5, 3.15, 10.9, 3.15, color=EMERALD, weight=3)
pill(s, 5.5, 2.6, "🔒  HTTPS / TLS 1.3  ·  AES-256", fill=EMERALD, color=DARK, size=11, w=2.4)
pill(s, 5.5, 3.3, "Event thumbnails + JSON metadata", fill=CARD, color=INK, size=10, w=2.4)

# bottom: compliance badges
rect(s, 0.6, 4.6, 12.2, 2.3, fill=CARD, line=BORDER)
txt(s, 0.8, 4.75, 12, 0.4, "COMPLIANCE & CONTROLS", size=11, color=DANGER, bold=True)
badges = [
    ("DPDP", "India 2023", SAFFRON),
    ("PDPA", "Singapore 2012", VIOLET),
    ("PDPA", "Malaysia 2010", EMERALD),
    ("GDPR", "EU-aligned", INK),
    ("PCI", "SAQ-A (Razorpay)", DANGER),
]
for i, (name, sub, c) in enumerate(badges):
    x = 0.9 + i * 2.4
    rect(s, x, 5.2, 2.15, 0.9, fill=PAPER, line=c)
    txt(s, x, 5.28, 2.15, 0.4, name, size=14, color=c, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x, 5.65, 2.15, 0.35, sub, size=10, color=MUTED, align=PP_ALIGN.CENTER)

controls = "JWT · bcrypt(cost=12) · HSTS · Rate limiting · OWASP Top 10 mitigated · Annual pentest · 99.5% uptime SLA"
txt(s, 0.8, 6.3, 12, 0.4, controls, size=10, color=INK,
    align=PP_ALIGN.CENTER, italic=True)
footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — COMPETITIVE EDGE
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "COMPETITIVE POSITIONING", size=11, color=SAFFRON, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "Why Vantag wins against legacy & big-brand",
    size=32, color=INK, bold=True)

# Table
headers = [("", 0.8), ("VANTAG", 4.2), ("Legacy AI-CCTV", 7.7), ("Consumer Cam", 10.8)]
rect(s, 0.6, 1.9, 12.2, 0.6, fill=DARK, rounded=True)
for label, x in headers:
    color = SAFFRON if label == "VANTAG" else CARD
    txt(s, x, 2.0, 3.5, 0.4, label, size=13, color=color, bold=True)

rows = [
    ("Cost per store",       "₹1.9K–9K /mo",   "₹4L–40L one-time", "₹8K+ DIY"),
    ("AI detections",         "12 types",       "3–6 (vendor lock)", "Motion only"),
    ("Setup time",            "< 30 minutes",   "2–3 weeks install", "1 hour"),
    ("Video privacy",         "Stays on-prem",  "Cloud-stored",      "Cloud-stored"),
    ("Hardware freedom",      "Any IP camera",  "Proprietary only",  "Locked brand"),
    ("Local language",        "11 languages",   "English only",      "English only"),
    ("Per-zone rules",        "Unlimited",      "Paid upgrade",      "None"),
]
for i, row in enumerate(rows):
    y = 2.55 + i * 0.58
    bg = CARD if i % 2 == 0 else PAPER
    rect(s, 0.6, y, 12.2, 0.56, fill=bg, rounded=False, line=BORDER)
    # Vantag column highlighted
    rect(s, 3.9, y, 3.3, 0.56, fill=EMERALD, rounded=False)
    for (col_text, col_x, col_w) in [(row[0], 0.75, 3.1), (row[1], 4.05, 3.1),
                                       (row[2], 7.55, 3.1), (row[3], 10.7, 2.5)]:
        tc = CARD if col_x == 4.05 else INK
        bold = col_x == 4.05 or col_x == 0.75
        txt(s, col_x, y + 0.11, col_w, 0.4, col_text, size=12, color=tc, bold=bold)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — TRACTION (placeholder)
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "GO-TO-MARKET", size=11, color=EMERALD, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "How we acquire 1,500 paying tenants",
    size=32, color=INK, bold=True)

channels = [
    ("🔍", "SEO + CONTENT", "Rank for 'AI CCTV for shop',\n'retail theft', 'DIY security'",
     "CAC target: ₹ 400"),
    ("📢", "GOOGLE ADS", "High-intent searches in\nall 3 regions",
     "CAC target: ₹ 1,200"),
    ("🏢", "MALL DIRECT SALES", "24 SG malls × 200 cams =\n4,800 potential cameras",
     "Enterprise deals"),
    ("🤝", "INSTALLER CHANNEL", "India CCTV installers get\n20% recurring commission",
     "Scale lever"),
    ("🔄", "POS PARTNERSHIPS", "Square, Pine Labs integrations\nfor cross-sell",
     "Free leads"),
    ("📱", "REFERRAL PROGRAM", "1 month free per successful\nreferral on both sides",
     "Viral coefficient"),
]
for i, (icon, title, desc, metric) in enumerate(channels):
    row, col = divmod(i, 3)
    x = 0.6 + col * 4.15
    y = 2.1 + row * 2.3
    rect(s, x, y, 3.95, 2.1, fill=CARD, line=BORDER)
    circle(s, x + 0.3, y + 0.25, 0.8, fill=EMERALD if i % 2 == 0 else VIOLET)
    txt(s, x + 0.3, y + 0.25, 0.8, 0.8, icon, size=26, color=CARD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 1.25, y + 0.25, 2.5, 0.4, title, size=13, color=INK, bold=True)
    txt(s, x + 1.25, y + 0.65, 2.6, 0.9, desc, size=11, color=MUTED)
    pill(s, x + 0.3, y + 1.55, metric, fill=DARK, color=CARD, size=9, w=3.4)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — ROADMAP
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "ROADMAP", size=11, color=VIOLET, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "From v1.0 GA to enterprise scale",
    size=32, color=INK, bold=True)

# Horizontal timeline
rect(s, 0.8, 4.0, 11.8, 0.05, fill=MUTED, rounded=False)

phases = [
    ("v0.9",  "BETA", "3 detectors\n1 region (SG)\n50 tenants", "Completed", EMERALD),
    ("v1.0",  "GA NOW", "12 detectors\n3 regions · Razorpay\nSuper-admin", "Current", VIOLET),
    ("v1.1",  "Q3 2026", "Deep POS integrations\nLoss-prevention pack\nJetson optimisation", "Next", SAFFRON),
    ("v1.2",  "Q4 2026", "Cloud-hybrid AI\nWhite-label\nCRM webhooks", "Planned", INK),
    ("v2.0",  "Q2 2027", "Enterprise SSO\nOn-prem option\nMulti-site federation", "Future", DANGER),
]
for i, (ver, date, items, status, color) in enumerate(phases):
    x = 0.8 + i * 2.45
    # marker
    circle(s, x + 0.95, 3.75, 0.55, fill=color)
    txt(s, x + 0.95, 3.75, 0.55, 0.55, ver, size=11, color=CARD, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # card
    above = i % 2 == 0
    cy = 1.9 if above else 4.5
    rect(s, x, cy, 2.25, 1.65, fill=CARD, line=color)
    txt(s, x + 0.1, cy + 0.1, 2.1, 0.35, date, size=11, color=color, bold=True)
    txt(s, x + 0.1, cy + 0.45, 2.1, 1.0, items, size=10, color=INK)
    pill(s, x + 0.1, cy + 1.3, status, fill=color, color=CARD, size=8, w=1.0)
    # connector
    line(s, x + 1.075, cy + 1.65 if above else cy, x + 1.075, 3.75, color=color, weight=1.5)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — KPIs YEAR 1
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "YEAR-1 TARGETS", size=11, color=SAFFRON, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "What success looks like",
    size=32, color=INK, bold=True)

kpis = [
    ("1,500",   "paying tenants",        VIOLET),
    ("USD 45K", "MRR exit run-rate",     EMERALD),
    ("68%",     "gross margin",          SAFFRON),
    ("< 6 mo",  "CAC payback",           VIOLET),
    ("> 92%",   "AI detection accuracy", EMERALD),
    ("< 8%",    "false-positive rate",   DANGER),
    ("> 99.5%", "cloud uptime",          INK),
    ("> 45",    "net promoter score",    SAFFRON),
]
for i, (big, lbl, color) in enumerate(kpis):
    row, col = divmod(i, 4)
    x = 0.6 + col * 3.1
    y = 2.2 + row * 2.3
    rect(s, x, y, 2.95, 2.1, fill=CARD, line=BORDER)
    rect(s, x, y, 2.95, 0.15, fill=color, rounded=False)
    txt(s, x + 0.1, y + 0.35, 2.75, 1.1, big, size=44, color=color, bold=True,
        align=PP_ALIGN.CENTER)
    txt(s, x + 0.1, y + 1.5, 2.75, 0.5, lbl, size=13, color=INK,
        align=PP_ALIGN.CENTER)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 — WHO IT'S FOR
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "TARGET VERTICALS", size=11, color=EMERALD, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "Built for more than just retail",
    size=32, color=INK, bold=True)

verticals = [
    ("🏪", "RETAIL SHOPS",    "2–10 cameras",  "Theft · Shrinkage · Staff"),
    ("🏬", "SHOPPING MALLS",  "16–60 cameras", "Crowd · Fall · Intrusion"),
    ("🏥", "HOSPITALS",       "20–50 cameras", "Patient fall · Restricted zones"),
    ("💊", "CLINICS / PHARMA","3–12 cameras",  "After-hours · Inventory"),
    ("📮", "POST OFFICES",    "4–10 cameras",  "Public safety · Tamper"),
    ("🚔", "POLICE STATIONS", "6–20 cameras",  "Custody · Restricted access"),
    ("🏢", "SMALL OFFICES",   "3–15 cameras",  "After-hours · Tailgating"),
    ("🏠", "HOMES",           "2–8 cameras",   "Elderly fall · Intrusion"),
]
for i, (icon, title, cams, uses) in enumerate(verticals):
    row, col = divmod(i, 4)
    x = 0.6 + col * 3.1
    y = 2.1 + row * 2.35
    rect(s, x, y, 2.95, 2.15, fill=CARD, line=BORDER)
    rect(s, x, y, 2.95, 0.8, fill=VIOLET, rounded=False)
    txt(s, x, y, 2.95, 0.8, icon, size=36, color=CARD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 0.15, y + 0.9, 2.65, 0.4, title, size=13, color=INK, bold=True)
    txt(s, x + 0.15, y + 1.3, 2.65, 0.3, cams, size=10, color=EMERALD, bold=True)
    txt(s, x + 0.15, y + 1.6, 2.65, 0.5, uses, size=10, color=MUTED)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17 — DOCUMENTATION / CREDIBILITY
# ═══════════════════════════════════════════════════════════════════════
s = add_slide()
txt(s, 0.6, 0.35, 12, 0.5, "BUILT RIGHT", size=11, color=VIOLET, bold=True)
txt(s, 0.6, 0.75, 12, 0.9, "Enterprise-grade documentation & governance",
    size=32, color=INK, bold=True)

# Left: big stat
rect(s, 0.6, 2.0, 5.0, 4.9, fill=DARK, rounded=True)
txt(s, 0.6, 2.5, 5.0, 1.8, "22", size=140, color=SAFFRON, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, 0.6, 4.5, 5.0, 0.5, "PRODUCTION DOCUMENTS", size=14, color=CARD, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, 0.6, 5.0, 5.0, 0.4,
    "PRD · BRD · SRS · HLD · LLD",
    size=12, color=CARD, align=PP_ALIGN.CENTER)
txt(s, 0.6, 5.4, 5.0, 0.4,
    "API · DB · Security · Dev · ML",
    size=12, color=CARD, align=PP_ALIGN.CENTER)
txt(s, 0.6, 5.8, 5.0, 0.4,
    "Ops · QA · Compliance · User Guides",
    size=12, color=CARD, align=PP_ALIGN.CENTER)
pill(s, 2.3, 6.35, "On GitHub · MIT license", fill=EMERALD, color=CARD, size=10, w=2.1)

# Right: practices
rect(s, 5.9, 2.0, 6.9, 4.9, fill=CARD, line=BORDER)
txt(s, 6.1, 2.15, 6.5, 0.4, "ENGINEERING PRACTICES", size=11, color=EMERALD, bold=True)
practices = [
    ("✓ GitHub CI/CD",             "Every push tested & deployed"),
    ("✓ 95%+ test pass rate",       "62 cases across BE/FE/AI"),
    ("✓ OWASP Top-10 mitigated",    "Annual external pentest"),
    ("✓ 99.5% uptime SLA",          "Rolling deployments"),
    ("✓ Structured logs + Sentry",  "Full observability"),
    ("✓ Daily DB backups",          "Cross-region · 30-day retention"),
    ("✓ Multi-tenant isolation",    "Row-level security in PostgreSQL"),
    ("✓ Edge-first architecture",   "No customer video in cloud"),
]
for i, (chk, sub) in enumerate(practices):
    y = 2.7 + i * 0.5
    txt(s, 6.2, y, 3.2, 0.4, chk, size=12, color=EMERALD, bold=True)
    txt(s, 9.2, y, 3.5, 0.4, sub, size=11, color=INK)

footer(s, nextn(), TOTAL)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18 — CTA / THANK YOU
# ═══════════════════════════════════════════════════════════════════════
s = add_slide(bg=DARK)
rect(s, 0, 0, 13.333, 0.15, fill=VIOLET, rounded=False)
rect(s, 4.44, 0, 4.44, 0.15, fill=SAFFRON, rounded=False)
rect(s, 8.88, 0, 4.44, 0.15, fill=EMERALD, rounded=False)

txt(s, 0.6, 1.8, 12.2, 0.7, "PARTNER WITH US", size=14, color=SAFFRON, bold=True,
    align=PP_ALIGN.CENTER)
txt(s, 0.6, 2.35, 12.2, 1.2, "Protect every retail\nbusiness on the planet.",
    size=48, color=CARD, bold=True, align=PP_ALIGN.CENTER)

# CTA cards
ctas = [
    ("🚀",  "TRY IT",  "retail-vantag.com\nretailnazar.com · jagajaga.my", VIOLET),
    ("💼",  "PARTNER", "Installers · POS vendors\nMall chains · Distributors", SAFFRON),
    ("💰",  "INVEST",  "Seed round opening\nanandsg.kumar@gmail.com", EMERALD),
]
for i, (icon, title, desc, color) in enumerate(ctas):
    x = 0.8 + i * 4.1
    rect(s, x, 4.8, 3.9, 1.9, fill=INK, line=color)
    circle(s, x + 0.3, 5.05, 0.7, fill=color)
    txt(s, x + 0.3, 5.05, 0.7, 0.7, icon, size=24, color=CARD,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + 1.2, 5.0, 2.6, 0.4, title, size=16, color=color, bold=True)
    txt(s, x + 1.2, 5.45, 2.6, 1.1, desc, size=11, color=CARD)

txt(s, 0.6, 6.9, 12.2, 0.3,
    "support@retail-vantag.com   ·   github.com/anandindiakr/Vantag",
    size=12, color=MUTED, align=PP_ALIGN.CENTER)

# ─── Save ──────────────────────────────────────────────────────────────
import os
out = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..",
                   "Vantag_Platform_Overview.pptx")
out = os.path.abspath(out)
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
